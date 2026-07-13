#!/usr/bin/env python3
import argparse
import json
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


SLIDE_WIDTH = 13.333
SLIDE_HEIGHT = 7.5


THEMES = {
    "editorial": {
        "paper": "F5F2EA",
        "ink": "18251F",
        "accent": "0F766E",
        "muted": "66736C",
        "soft": "DCE8E2",
        "white": "FFFFFF",
        "title_font": "Georgia",
        "body_font": "Microsoft YaHei",
    },
    "swiss": {
        "paper": "F5F5F1",
        "ink": "111111",
        "accent": "155EEF",
        "muted": "606060",
        "soft": "E5E5DF",
        "white": "FFFFFF",
        "title_font": "Arial",
        "body_font": "Microsoft YaHei",
    },
}


def color(value):
    value = value.lstrip("#")
    return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))


def clean_text(value, fallback=""):
    if value is None:
        return fallback
    return str(value).strip() or fallback


def title_size(text, maximum=38, minimum=24):
    length = max(1, len(clean_text(text)))
    if length <= 10:
        return maximum
    if length <= 18:
        return max(minimum, maximum - 5)
    if length <= 28:
        return max(minimum, maximum - 10)
    return minimum


def set_background(slide, hex_color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color(hex_color)


def add_rect(slide, x, y, width, height, fill, line=None):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color(fill)
    if line:
        shape.line.color.rgb = color(line)
    else:
        shape.line.fill.background()
    return shape


def add_text(
    slide,
    text,
    x,
    y,
    width,
    height,
    size,
    fill,
    font,
    bold=False,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
    margin=0,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(width), Inches(height))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(margin)
    frame.margin_right = Inches(margin)
    frame.margin_top = Inches(margin)
    frame.margin_bottom = Inches(margin)
    frame.vertical_anchor = valign
    paragraph = frame.paragraphs[0]
    paragraph.text = clean_text(text)
    paragraph.alignment = align
    paragraph.line_spacing = 1.05
    if paragraph.runs:
        run = paragraph.runs[0]
        run.font.name = font
        run.font.size = Pt(size)
        run.font.color.rgb = color(fill)
        run.font.bold = bold
    return box


def add_bullet_list(slide, values, x, y, width, height, theme, size=18, limit=5):
    items = [clean_text(item) for item in values if clean_text(item)][:limit]
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(width), Inches(height))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = 0
    frame.margin_right = 0
    frame.margin_top = 0
    frame.margin_bottom = 0
    for index, item in enumerate(items):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = "• " + item
        paragraph.space_after = Pt(14)
        paragraph.line_spacing = 1.12
        if paragraph.runs:
            run = paragraph.runs[0]
            run.font.name = theme["body_font"]
            run.font.size = Pt(size)
            run.font.color.rgb = color(theme["ink"])
    return box


def add_footer(slide, page, total, theme, inverse=False):
    ink = theme["paper"] if inverse else theme["muted"]
    rule = theme["accent"] if inverse else theme["ink"]
    add_rect(slide, 0.65, 7.03, 11.55, 0.018, rule)
    add_text(
        slide,
        "AEGISY",
        0.68,
        7.09,
        1.5,
        0.2,
        8,
        ink,
        "Arial",
        bold=True,
    )
    add_text(
        slide,
        f"{page:02d} / {total:02d}",
        11.25,
        7.07,
        1.0,
        0.22,
        9,
        ink,
        "Arial",
        align=PP_ALIGN.RIGHT,
    )


def add_header(slide, item, page, total, theme):
    kicker = clean_text(item.get("kicker"), f"SECTION {page:02d}")
    add_text(slide, kicker.upper(), 0.72, 0.48, 3.8, 0.28, 10, theme["accent"], "Arial", True)
    title = clean_text(item.get("title"), f"第 {page} 部分")
    add_text(
        slide,
        title,
        0.72,
        0.86,
        11.75,
        0.75,
        title_size(title, 31, 22),
        theme["ink"],
        theme["title_font"],
        True,
    )
    add_footer(slide, page, total, theme)


def add_notes(slide, notes):
    notes = clean_text(notes)
    if not notes:
        return
    try:
        slide.notes_slide.notes_text_frame.text = notes
    except (AttributeError, NotImplementedError):
        pass


def render_cover(presentation, item, page, total, theme_name, theme):
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    title = clean_text(item.get("title"), "演示文稿")
    subtitle = clean_text(item.get("subtitle"))
    kicker = clean_text(item.get("kicker"), "AEGISY PRESENTATION")
    if theme_name == "swiss":
        set_background(slide, theme["paper"])
        add_rect(slide, 0, 0, 2.75, SLIDE_HEIGHT, theme["accent"])
        add_text(slide, f"{page:02d}", 0.45, 0.45, 1.8, 0.8, 46, theme["white"], "Arial")
        add_text(slide, kicker.upper(), 3.25, 0.8, 4.8, 0.3, 10, theme["accent"], "Arial", True)
        add_text(
            slide,
            title,
            3.22,
            1.55,
            9.2,
            2.3,
            title_size(title, 46, 30),
            theme["ink"],
            theme["title_font"],
            False,
            valign=MSO_ANCHOR.MIDDLE,
        )
        add_rect(slide, 3.25, 4.45, 2.0, 0.08, theme["accent"])
        add_text(slide, subtitle, 3.25, 4.8, 8.3, 0.9, 18, theme["muted"], theme["body_font"])
        add_footer(slide, page, total, theme)
    else:
        set_background(slide, theme["ink"])
        add_rect(slide, 0.78, 0.8, 0.16, 5.35, theme["accent"])
        add_text(slide, kicker.upper(), 1.35, 1.02, 5.5, 0.35, 10, theme["soft"], "Arial", True)
        add_text(
            slide,
            title,
            1.3,
            1.72,
            10.65,
            2.4,
            title_size(title, 42, 29),
            theme["paper"],
            theme["title_font"],
            True,
            valign=MSO_ANCHOR.MIDDLE,
        )
        add_text(slide, subtitle, 1.35, 4.55, 9.3, 0.8, 18, theme["soft"], theme["body_font"])
        add_footer(slide, page, total, theme, inverse=True)
    add_notes(slide, item.get("notes"))


def render_section(presentation, item, page, total, theme):
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    set_background(slide, theme["accent"])
    add_text(slide, f"{page:02d}", 0.72, 0.62, 2.0, 0.8, 48, theme["white"], "Arial")
    add_text(
        slide,
        clean_text(item.get("title"), "章节"),
        2.9,
        1.65,
        9.25,
        2.0,
        title_size(item.get("title"), 42, 28),
        theme["white"],
        theme["title_font"],
        True,
        valign=MSO_ANCHOR.MIDDLE,
    )
    add_rect(slide, 2.92, 4.08, 2.1, 0.08, theme["white"])
    add_text(slide, item.get("subtitle"), 2.92, 4.45, 8.5, 1.0, 18, theme["white"], theme["body_font"])
    add_footer(slide, page, total, theme, inverse=True)
    add_notes(slide, item.get("notes"))


def render_statement(presentation, item, page, total, theme):
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    set_background(slide, theme["paper"])
    add_rect(slide, 0.72, 0.58, 1.8, 0.08, theme["accent"])
    add_text(slide, clean_text(item.get("kicker"), "KEY IDEA"), 0.72, 0.8, 3.5, 0.3, 10, theme["accent"], "Arial", True)
    statement = clean_text(item.get("quote"), clean_text(item.get("title"), "核心观点"))
    add_text(
        slide,
        statement,
        0.78,
        1.45,
        11.75,
        3.55,
        title_size(statement, 36, 25),
        theme["ink"],
        theme["title_font"],
        True,
        valign=MSO_ANCHOR.MIDDLE,
    )
    add_text(slide, item.get("subtitle"), 7.0, 5.4, 5.45, 0.7, 16, theme["muted"], theme["body_font"], align=PP_ALIGN.RIGHT)
    add_footer(slide, page, total, theme)
    add_notes(slide, item.get("notes"))


def render_bullets(presentation, item, page, total, theme):
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    set_background(slide, theme["paper"])
    add_header(slide, item, page, total, theme)
    add_bullet_list(slide, item.get("bullets") or [], 0.86, 1.95, 8.15, 4.55, theme, 18)
    add_rect(slide, 9.6, 1.95, 0.08, 4.25, theme["accent"])
    add_text(slide, f"{page:02d}", 10.05, 2.0, 2.0, 0.85, 42, theme["accent"], "Arial")
    add_text(slide, item.get("subtitle"), 10.05, 3.05, 2.2, 2.35, 15, theme["muted"], theme["body_font"])
    add_notes(slide, item.get("notes"))


def render_comparison(presentation, item, page, total, theme):
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    set_background(slide, theme["paper"])
    add_header(slide, item, page, total, theme)
    columns = item.get("columns") or []
    if len(columns) < 2:
        bullets = item.get("bullets") or []
        midpoint = max(1, len(bullets) // 2)
        columns = [
            {"title": "当前", "bullets": bullets[:midpoint]},
            {"title": "目标", "bullets": bullets[midpoint:]},
        ]
    add_rect(slide, 6.64, 1.95, 0.03, 4.35, theme["ink"])
    for index, column in enumerate(columns[:2]):
        x = 0.82 if index == 0 else 6.98
        add_text(slide, clean_text(column.get("title"), f"方案 {index + 1}"), x, 2.0, 5.25, 0.55, 22, theme["accent"], theme["title_font"], True)
        add_bullet_list(slide, column.get("bullets") or [], x, 2.82, 5.15, 3.25, theme, 16, 4)
    add_notes(slide, item.get("notes"))


def render_process(presentation, item, page, total, theme):
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    set_background(slide, theme["paper"])
    add_header(slide, item, page, total, theme)
    steps = item.get("steps") or item.get("bullets") or []
    normalized = []
    for value in steps[:4]:
        if isinstance(value, dict):
            normalized.append(value)
        else:
            text = clean_text(value)
            parts = text.replace(":", "：").split("：", 1)
            normalized.append({"title": parts[0], "description": parts[1] if len(parts) > 1 else ""})
    count = max(1, len(normalized))
    available = 11.85
    gap = 0.22
    width = (available - gap * (count - 1)) / count
    for index, step in enumerate(normalized):
        x = 0.72 + index * (width + gap)
        add_rect(slide, x, 2.12, width, 0.08, theme["accent"])
        add_text(slide, f"{index + 1:02d}", x, 2.45, width, 0.65, 30, theme["accent"], "Arial")
        add_text(slide, clean_text(step.get("title"), f"步骤 {index + 1}"), x, 3.2, width, 0.75, 18, theme["ink"], theme["body_font"], True)
        add_text(slide, step.get("description"), x, 4.15, width, 1.45, 13, theme["muted"], theme["body_font"])
    add_notes(slide, item.get("notes"))


def render_metrics(presentation, item, page, total, theme):
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    set_background(slide, theme["paper"])
    add_header(slide, item, page, total, theme)
    metrics = item.get("metrics") or []
    if not metrics:
        metrics = []
        for bullet in (item.get("bullets") or [])[:4]:
            parts = clean_text(bullet).replace(":", "：").split("：", 1)
            metrics.append({"value": parts[0], "label": parts[1] if len(parts) > 1 else ""})
    metrics = metrics[:4]
    count = max(1, len(metrics))
    width = 11.7 / count
    for index, metric in enumerate(metrics):
        x = 0.75 + index * width
        if index:
            add_rect(slide, x - 0.12, 2.2, 0.025, 3.65, theme["soft"])
        add_text(slide, clean_text(metric.get("value"), "—"), x, 2.25, width - 0.3, 1.25, 34, theme["accent"], "Arial")
        add_text(slide, metric.get("label"), x, 3.7, width - 0.3, 0.75, 17, theme["ink"], theme["body_font"], True)
        add_text(slide, metric.get("note"), x, 4.65, width - 0.3, 1.0, 12, theme["muted"], theme["body_font"])
    add_notes(slide, item.get("notes"))


def render_closing(presentation, item, page, total, theme):
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    set_background(slide, theme["ink"])
    add_text(slide, clean_text(item.get("kicker"), "TAKEAWAY"), 0.82, 0.8, 3.4, 0.3, 10, theme["accent"], "Arial", True)
    title = clean_text(item.get("title"), "谢谢")
    add_text(slide, title, 0.78, 1.55, 11.4, 1.7, title_size(title, 42, 28), theme["paper"], theme["title_font"], True, valign=MSO_ANCHOR.MIDDLE)
    add_rect(slide, 0.82, 3.72, 2.15, 0.08, theme["accent"])
    add_text(slide, clean_text(item.get("quote"), item.get("subtitle")), 0.82, 4.15, 9.8, 1.25, 20, theme["soft"], theme["body_font"])
    add_footer(slide, page, total, theme, inverse=True)
    add_notes(slide, item.get("notes"))


RENDERERS = {
    "cover": render_cover,
    "section": render_section,
    "statement": render_statement,
    "bullets": render_bullets,
    "comparison": render_comparison,
    "process": render_process,
    "metrics": render_metrics,
    "closing": render_closing,
}


def build(spec, output):
    presentation = Presentation()
    presentation.slide_width = Inches(SLIDE_WIDTH)
    presentation.slide_height = Inches(SLIDE_HEIGHT)

    theme_name = clean_text(spec.get("theme"), "editorial").lower()
    if theme_name not in THEMES:
        theme_name = "editorial"
    theme = THEMES[theme_name]
    source_slides = spec.get("slides") if isinstance(spec.get("slides"), list) else []
    prepared = [item for item in source_slides[:20] if isinstance(item, dict)]
    if not prepared or clean_text(prepared[0].get("layout")).lower() != "cover":
        prepared.insert(
            0,
            {
                "layout": "cover",
                "title": clean_text(spec.get("title"), "演示文稿"),
                "subtitle": clean_text(spec.get("subtitle"), "由 Aegisy Skills 生成"),
            },
        )
    else:
        prepared[0].setdefault("title", clean_text(spec.get("title"), "演示文稿"))
        prepared[0].setdefault("subtitle", clean_text(spec.get("subtitle")))

    total = len(prepared)
    for page, item in enumerate(prepared, start=1):
        layout = clean_text(item.get("layout"), "bullets").lower()
        renderer = RENDERERS.get(layout, render_bullets)
        if layout == "cover":
            renderer(presentation, item, page, total, theme_name, theme)
        else:
            renderer(presentation, item, page, total, theme)

    output_path = Path(output).expanduser().resolve()
    output_path.parent.mkdir(parents=True, exist_ok=True)
    presentation.save(output_path)
    print(
        json.dumps(
            {"ok": True, "path": str(output_path), "slides": len(prepared), "theme": theme_name},
            ensure_ascii=False,
        )
    )


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--spec", required=True)
    parser.add_argument("--out", required=True)
    args = parser.parse_args()
    with open(args.spec, "r", encoding="utf-8") as source:
        build(json.load(source), args.out)


if __name__ == "__main__":
    main()
